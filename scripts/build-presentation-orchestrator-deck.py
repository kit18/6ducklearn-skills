from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "out"
PPTX_PATH = OUT_DIR / "presentation-orchestrator-how-it-works.pptx"
MD_PATH = OUT_DIR / "presentation-orchestrator-how-it-works.md"


def rgb(value):
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


COLORS = {
    "white": rgb("#ffffff"),
    "bg": rgb("#f6f7f9"),
    "ink": rgb("#111827"),
    "muted": rgb("#5b6472"),
    "line": rgb("#b9c0cc"),
    "rule": rgb("#d8dde6"),
    "navy": rgb("#172033"),
    "teal": rgb("#00736b"),
    "teal_light": rgb("#e0f2ef"),
    "blue": rgb("#2446a7"),
    "blue_light": rgb("#e8edff"),
    "amber": rgb("#b46a00"),
    "amber_light": rgb("#fff1cf"),
    "red": rgb("#bd2445"),
    "red_light": rgb("#fde3e9"),
    "green": rgb("#16803a"),
    "green_light": rgb("#e3f5e8"),
    "gray": rgb("#eef1f5"),
    "dark_gray": rgb("#374151"),
}


def set_bg(slide, color=COLORS["bg"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_run(run, size=10, bold=False, color=None):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]


def text(slide, x, y, w, h, value, size=10, bold=False, color=None,
         align=None, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = value
    if align:
        p.alignment = align
    set_run(p.runs[0], size=size, bold=bold, color=color)
    return box


def rect(slide, x, y, w, h, fill=COLORS["white"], line=COLORS["line"], width=0.7):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(width)
    return s


def hline(slide, x1, y, x2, color=COLORS["rule"], width=0.8):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y),
                                   Inches(x2), Inches(y))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def connector(slide, x1, y1, x2, y2, color=COLORS["line"], width=1.0):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def header(slide, section, title, source=None):
    text(slide, 0.55, 0.22, 2.8, 0.2, section.upper(), size=7.2,
         bold=True, color=COLORS["teal"])
    text(slide, 0.55, 0.48, 11.2, 0.62, title, size=19.2, bold=True,
         color=COLORS["ink"])
    hline(slide, 0.55, 1.15, 12.75, color=COLORS["rule"], width=0.7)
    if source:
        text(slide, 0.55, 7.08, 10.4, 0.2, f"Source: {source}",
             size=7.0, color=COLORS["muted"])
    text(slide, 11.5, 7.08, 1.25, 0.2, "6ducklearn", size=7.0,
         bold=True, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def note(slide, note_text):
    try:
        slide.notes_slide.notes_text_frame.text = note_text
    except Exception:
        pass


def label(slide, x, y, w, h, title, body=None, fill=COLORS["white"],
          line=COLORS["line"], title_color=None, title_size=9.2, body_size=7.4,
          align=PP_ALIGN.CENTER):
    rect(slide, x, y, w, h, fill=fill, line=line)
    text(slide, x + 0.06, y + 0.05, w - 0.12, 0.22, title, size=title_size,
         bold=True, color=title_color or COLORS["ink"], align=align)
    if body:
        text(slide, x + 0.08, y + 0.34, w - 0.16, h - 0.42, body,
             size=body_size, color=COLORS["muted"], align=align)


def callout(slide, x, y, w, h, title, body, accent=COLORS["teal"],
            fill=COLORS["white"]):
    rect(slide, x, y, w, h, fill=fill, line=COLORS["line"])
    rect(slide, x, y, 0.06, h, fill=accent, line=accent)
    if h <= 0.72:
        text(slide, x + 0.16, y + 0.16, 2.3, h - 0.18, title, size=8.8,
             bold=True)
        text(slide, x + 2.25, y + 0.16, w - 2.45, h - 0.18, body, size=7.5,
             color=COLORS["muted"])
    else:
        text(slide, x + 0.16, y + 0.12, w - 0.24, 0.26, title, size=10.2,
             bold=True)
        text(slide, x + 0.16, y + 0.46, w - 0.24, h - 0.52, body, size=8.4,
             color=COLORS["muted"])


def table(slide, x, y, col_widths, row_h, headers, rows, header_fill=COLORS["navy"],
          header_text=COLORS["white"], font_size=7.3):
    for c, value in enumerate(headers):
        cx = x + sum(col_widths[:c])
        rect(slide, cx, y, col_widths[c], 0.35, fill=header_fill,
             line=header_fill)
        text(slide, cx + 0.05, y + 0.09, col_widths[c] - 0.1, 0.14,
             value, size=font_size, bold=True, color=header_text,
             align=PP_ALIGN.CENTER)
    y0 = y + 0.35
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cx = x + sum(col_widths[:c])
            fill = COLORS["white"] if r % 2 == 0 else COLORS["gray"]
            rect(slide, cx, y0 + r * row_h, col_widths[c], row_h,
                 fill=fill, line=COLORS["line"], width=0.5)
            text(slide, cx + 0.06, y0 + r * row_h + 0.06,
                 col_widths[c] - 0.12, row_h - 0.08, str(value),
                 size=font_size, bold=(c == 0), color=COLORS["ink"] if c == 0 else COLORS["muted"])


def score_cell(slide, x, y, label_text, status, detail, color):
    rect(slide, x, y, 2.85, 0.92, fill=COLORS["white"], line=COLORS["line"])
    rect(slide, x, y, 0.08, 0.92, fill=color, line=color)
    text(slide, x + 0.18, y + 0.1, 1.5, 0.16, label_text, size=8.5, bold=True)
    text(slide, x + 1.92, y + 0.1, 0.7, 0.16, status, size=8.0, bold=True,
         color=color, align=PP_ALIGN.RIGHT)
    text(slide, x + 0.18, y + 0.38, 2.45, 0.42, detail, size=7.2,
         color=COLORS["muted"])


def tag(slide, x, y, w, value, fill, color=None):
    rect(slide, x, y, w, 0.28, fill=fill, line=fill)
    text(slide, x + 0.05, y + 0.07, w - 0.1, 0.1, value.upper(), size=6.6,
         bold=True, color=color or COLORS["white"], align=PP_ALIGN.CENTER)


def route_tile(slide, x, y, w, h, method, decision, reason, color, selected=False):
    fill = COLORS["white"] if not selected else COLORS["teal_light"]
    rect(slide, x, y, w, h, fill=fill, line=color, width=1.0 if selected else 0.7)
    rect(slide, x, y, 0.08, h, fill=color, line=color)
    tag(slide, x + w - 1.25, y + 0.12, 1.02, decision, color)
    text(slide, x + 0.22, y + 0.14, w - 1.62, 0.22, method, size=9.3,
         bold=True, color=COLORS["ink"])
    text(slide, x + 0.22, y + 0.52, w - 0.42, h - 0.64, reason, size=7.4,
         color=COLORS["muted"])


def storyboard_card(slide, x, y, w, h, number, title, exhibit, proof):
    rect(slide, x, y, w, h, fill=COLORS["white"], line=COLORS["line"])
    rect(slide, x, y, 0.42, h, fill=COLORS["teal"], line=COLORS["teal"])
    text(slide, x + 0.06, y + 0.18, 0.3, 0.2, str(number), size=10.0,
         bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    text(slide, x + 0.58, y + 0.12, w - 0.78, 0.34, title, size=8.7,
         bold=True, color=COLORS["ink"])
    text(slide, x + 0.58, y + 0.58, w - 0.78, 0.2, exhibit, size=7.3,
         bold=True, color=COLORS["teal"])
    text(slide, x + 0.58, y + 0.84, w - 0.78, h - 0.96, proof, size=7.0,
         color=COLORS["muted"])


def delta_row(slide, y, gate, before, after, color):
    x0 = 0.78
    rect(slide, x0, y, 2.15, 0.52, fill=COLORS["white"], line=COLORS["line"])
    text(slide, x0 + 0.15, y + 0.15, 1.85, 0.16, gate, size=8.4,
         bold=True, color=COLORS["ink"])
    rect(slide, 3.22, y, 3.55, 0.52, fill=COLORS["gray"], line=COLORS["line"])
    text(slide, 3.38, y + 0.12, 3.2, 0.2, before, size=7.3,
         color=COLORS["muted"])
    connector(slide, 6.92, y + 0.26, 7.38, y + 0.26, color=color, width=1.6)
    rect(slide, 7.55, y, 4.65, 0.52, fill=COLORS["white"], line=color)
    text(slide, 7.72, y + 0.12, 4.25, 0.2, after, size=7.3,
         bold=True, color=COLORS["ink"])


def validation_step(slide, x, y, number, title, proof, color):
    rect(slide, x, y, 2.1, 1.0, fill=COLORS["white"], line=COLORS["line"])
    rect(slide, x, y, 0.32, 1.0, fill=color, line=color)
    text(slide, x + 0.07, y + 0.38, 0.18, 0.12, str(number), size=8.0,
         bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    text(slide, x + 0.48, y + 0.13, 1.4, 0.18, title, size=8.3,
         bold=True, color=COLORS["ink"])
    text(slide, x + 0.48, y + 0.43, 1.42, 0.42, proof, size=6.8,
         color=COLORS["muted"])


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["white"])
    rect(slide, 0, 0, 13.333, 0.16, fill=COLORS["teal"], line=COLORS["teal"])
    text(slide, 0.72, 0.82, 7.2, 0.35, "Recommendation", size=9,
         bold=True, color=COLORS["teal"])
    text(slide, 0.72, 1.28, 8.1, 0.95,
         "Use Presentation Orchestrator when AI slide generation must become decision preparation",
         size=25, bold=True, color=COLORS["ink"])
    text(slide, 0.74, 2.48, 7.4, 0.52,
         "The skill is not a slide beautifier. It is a consulting-style harness for decks where a room must approve, choose, fund, stop, or assign.",
         size=14, color=COLORS["muted"])
    callout(slide, 8.7, 1.08, 3.85, 2.8, "Decision ask",
            "Pilot it on high-stakes business decks that need a clear POV, source-aware evidence, rejected paths, red-team rehearsal, and a specific next action.",
            accent=COLORS["teal"], fill=COLORS["teal_light"])
    table(slide, 0.78, 4.25, [2.35, 2.35, 2.35, 2.35, 2.35], 0.58,
          ["Use for", "Avoid for", "Core proof", "Human role", "Output"],
          [[
              "approval, strategy, GTM, roadmap",
              "quick cosmetic formatting",
              "case-led deck packet",
              "judge assumptions and risk",
              "PPTX/PDF + speaker notes",
          ]],
          header_fill=COLORS["navy"], font_size=7.2)
    text(slide, 0.78, 6.2, 11.6, 0.26,
         "Main claim: the skill improves the presenter's decision readiness, not merely slide aesthetics.",
         size=12.2, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)
    text(slide, 10.6, 6.9, 1.8, 0.2, "2026-06-26", size=7.2,
         color=COLORS["muted"], align=PP_ALIGN.RIGHT)
    note(slide, "Open with the recommendation and decision ask. Do not explain the file structure yet.")


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "executive summary",
           "The skill is worth using only when it can prove a stronger decision deck")
    text(slide, 0.75, 1.38, 3.3, 0.35, "Recommendation", size=9.2,
         bold=True, color=COLORS["teal"])
    text(slide, 0.75, 1.78, 3.25, 1.15,
         "Use the skill as a deck-strategy partner for high-stakes decisions; do not position it as a generic PPT generator.",
         size=15.0, bold=True)
    rows = [
        ("1. Decision spine", "Turns rough input into an issue question, route card, and explicit ask."),
        ("2. Evidence discipline", "Separates user facts, source-backed facts, assumptions, inference, and gaps."),
        ("3. Rehearsal loop", "Forces sponsor, skeptic, operator, and evidence-review objections before final output."),
        ("4. Verified artifact path", "This deck is generated reproducibly and validated as PPTX/PDF; package checks pass."),
    ]
    table(slide, 4.45, 1.38, [2.25, 5.25], 0.62, ["Proof point", "What it changes"],
          rows, header_fill=COLORS["teal"], font_size=8)
    callout(slide, 0.8, 5.2, 11.45, 0.9, "Consulting-team critique absorbed",
            "The prior version explained the process. This version uses one synthetic GTM case to prove judgment: selected method, rejected paths, evidence gaps, red-team changes, and readiness gates.",
            accent=COLORS["blue"], fill=COLORS["blue_light"])
    note(slide, "This slide is the answer-first executive summary. It should make the adoption logic clear in under two minutes.")


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "mini case",
           "A messy GTM request becomes a decision question the room can answer",
           "Illustrative synthetic case; no market facts claimed")
    callout(slide, 0.72, 1.42, 3.45, 3.7, "Raw user input",
            "We need a deck on whether to invest in a vertical GTM channel. We have notes, early customer signals, and leadership wants a recommendation next week.",
            accent=COLORS["red"], fill=COLORS["red_light"])
    connector(slide, 4.35, 3.28, 5.3, 3.28, color=COLORS["dark_gray"], width=1.6)
    label(slide, 5.45, 2.55, 2.55, 1.45, "Converted issue question",
          "Should we approve a gated vertical GTM pilot for the next 6 weeks, or defer until stronger evidence exists?",
          fill=COLORS["teal_light"], line=COLORS["teal"], title_size=9.4, body_size=8)
    connector(slide, 8.12, 3.28, 9.0, 3.28, color=COLORS["dark_gray"], width=1.6)
    callout(slide, 9.15, 1.42, 3.35, 3.7, "Decision-ready ask",
            "Approve a limited pilot with owner, proof gates, and kill criteria. Reject full-scale rollout until pipeline quality, CAC risk, and sales capacity are proven.",
            accent=COLORS["green"], fill=COLORS["green_light"])
    table(slide, 0.9, 5.48, [2.2, 2.5, 2.6, 2.6, 2.2], 0.45,
          ["Audience", "Decision", "Time box", "Evidence state", "Output"],
          [["growth leadership", "approve / defer", "6 weeks", "partial and mixed", "decision deck"]],
          header_fill=COLORS["navy"], font_size=7.1)
    note(slide, "This is the case used through the deck. Stress that it is synthetic and illustrates process, not market truth.")


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "route card",
           "The skill chooses the decision job before it chooses the framework",
           "methodology-router.md; presentation-diagram-source-catalog.md")
    callout(slide, 0.78, 1.35, 11.75, 0.58, "Decision job",
            "Approve, reject, or time-box a vertical GTM pilot under thin evidence. The room needs a decision, not a framework tour.",
            accent=COLORS["teal"], fill=COLORS["teal_light"])
    route_tile(slide, 0.78, 2.3, 5.55, 1.18, "Business Model Canvas",
               "reject", "Useful later for model design; wrong primary lens for a near-term pilot funding decision.",
               COLORS["red"])
    route_tile(slide, 6.78, 2.3, 5.55, 1.18, "Porter Five Forces",
               "reject", "Market structure may matter, but current evidence is too thin for an attractiveness diagnosis.",
               COLORS["red"])
    route_tile(slide, 0.78, 3.9, 5.55, 1.28, "Issue tree + option matrix",
               "select", "Compares scale, pilot, and defer against demand, economics, capability, and risk.",
               COLORS["teal"], selected=True)
    route_tile(slide, 6.78, 3.9, 5.55, 1.28, "Evidence ledger",
               "support", "Keeps user facts, assumptions, inference, and next checks visible before the ask.",
               COLORS["blue"])
    connector(slide, 3.55, 5.42, 3.55, 5.82, color=COLORS["teal"], width=1.3)
    connector(slide, 9.55, 5.42, 9.55, 5.82, color=COLORS["blue"], width=1.3)
    rect(slide, 1.2, 5.82, 10.7, 0.62, fill=COLORS["white"], line=COLORS["green"])
    text(slide, 1.42, 6.02, 10.25, 0.16,
         "Selected route: issue tree frames the decision; option matrix recommends the path; evidence ledger limits overclaiming.",
         size=8.8, bold=True, color=COLORS["ink"], align=PP_ALIGN.CENTER)
    note(slide, "This slide shows consultant judgment: two plausible frameworks are rejected with a reason.")


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "proof anchor",
           "The evidence ledger turns weak signals into explicit confidence and next checks",
           "deck-output-template.md; harness-loop-contract.md")
    rows = [
        ("Vertical segment shows interest", "user notes", "Medium", "sample biased toward warm prospects", "ask for pilot, not rollout"),
        ("Current team can execute pilot", "assumption", "Low", "owner capacity not confirmed", "name owner before approval"),
        ("Channel economics may work", "inference", "Low", "CAC/payback unproven", "define 2 evidence gates"),
        ("Full rollout is premature", "inference", "High", "proof is not strong enough", "reject scale-now option"),
        ("Defer loses learning", "assumption", "Medium", "opportunity cost unknown", "time-box the pilot"),
    ]
    table(slide, 0.62, 1.35, [2.35, 1.3, 1.0, 2.65, 3.25], 0.62,
          ["Claim", "Source type", "Conf.", "Gap / risk", "Slide implication"],
          rows, header_fill=COLORS["navy"], font_size=7.2)
    callout(slide, 0.9, 5.78, 11.25, 0.5, "What would change our mind",
            "If pilot CAC exceeds threshold, owner capacity is unavailable, or pipeline quality fails, the recommendation flips from pilot to defer.",
            accent=COLORS["red"], fill=COLORS["red_light"])
    note(slide, "This is the hero proof slide. Point out confidence and gaps before the recommendation.")


def slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "core exhibit",
           "The selected framework changes the decision: pilot beats scale-now and defer",
           "Illustrative synthetic case; no external market claim")
    # Issue tree on the left
    label(slide, 0.72, 2.82, 2.2, 0.74, "Issue question",
          "Fund a vertical GTM pilot?", fill=COLORS["teal_light"], line=COLORS["teal"])
    connector(slide, 2.92, 3.19, 3.35, 3.19, color=COLORS["teal"], width=1.3)
    connector(slide, 3.35, 1.76, 3.35, 5.08, color=COLORS["teal"], width=1.0)
    branches = [
        ("Demand", "Are signals real?", 3.72, 1.42, COLORS["blue"]),
        ("Economics", "Can CAC/payback work?", 3.72, 2.52, COLORS["amber"]),
        ("Capability", "Can team execute?", 3.72, 3.62, COLORS["green"]),
        ("Risk", "What would kill it?", 3.72, 4.72, COLORS["red"]),
    ]
    for title, body, x, y, color in branches:
        connector(slide, 3.35, y + 0.33, x, y + 0.33, color=color, width=1.1)
        label(slide, x, y, 2.0, 0.66, title, body, fill=COLORS["white"],
              line=color, title_size=8.4, body_size=6.6)
    text(slide, 6.25, 1.42, 4.8, 0.24, "Option readout", size=9.0,
         bold=True, color=COLORS["blue"])
    route_tile(slide, 6.25, 1.82, 5.25, 0.86, "Scale now",
               "reject", "High upside, high risk; evidence does not support rollout.",
               COLORS["red"])
    route_tile(slide, 6.25, 3.0, 5.25, 0.92, "Gated pilot",
               "recommend", "Medium upside and controlled risk; fastest path to proof.",
               COLORS["green"], selected=True)
    route_tile(slide, 6.25, 4.28, 5.25, 0.86, "Defer",
               "backup", "Lowest risk, but delays learning and leaves demand unresolved.",
               COLORS["amber"])
    callout(slide, 6.25, 5.72, 5.25, 0.56, "Decision logic",
            "Pilot wins because it buys evidence without buying a false conclusion.",
            accent=COLORS["green"], fill=COLORS["green_light"])
    note(slide, "This replaces the framework gallery. The framework changes the recommendation and rejects alternatives.")


def slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "storyboard",
           "The deck packet turns analysis into five action-title slides")
    cards = [
        ("Leadership should approve a gated pilot", "Decision brief", "ask, risk, owner"),
        ("The issue is learning under thin evidence", "Issue tree", "demand, economics, capability, risk"),
        ("Evidence supports testing, not scaling", "Evidence ledger", "confidence, gaps, flip criteria"),
        ("Pilot beats scale-now and defer", "Option matrix", "tradeoffs and decision logic"),
        ("Approval requires gates and kill criteria", "Execution close", "owner, gates, next meeting"),
    ]
    for i, (title, exhibit, proof) in enumerate(cards, 1):
        x = 0.6 + (i - 1) * 2.48
        storyboard_card(slide, x, 1.55, 2.22, 2.45, i, title, exhibit, proof)
        if i < len(cards):
            connector(slide, x + 2.24, 2.78, x + 2.42, 2.78, color=COLORS["teal"], width=1.0)
    callout(slide, 0.95, 4.85, 11.25, 0.82, "Presenter packet",
            "The output is not only slide titles. Each page carries a proof standard, speaker note, appendix cue, and known objection so the presenter can rehearse the decision path.",
            accent=COLORS["blue"], fill=COLORS["blue_light"])
    text(slide, 1.1, 6.08, 10.9, 0.28,
         "Quality bar: if a slide cannot change the decision, support the ask, or expose a risk, it moves to appendix or disappears.",
         size=10.4, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)
    note(slide, "Show this as the output of the skill, not just an outline.")


def slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "red-team revision",
           "The rehearsal loop improves the recommendation before the room reacts",
           "harness-loop-contract.md")
    rows = [
        ("CFO", "What evidence justifies funding?", "Funding ask changed from rollout to time-boxed pilot with kill criteria."),
        ("Operator", "Who owns Monday execution?", "Added owner confirmation as approval condition."),
        ("Sales leader", "Will this distract from current pipeline?", "Added capacity check and pilot scope limit."),
        ("Evidence reviewer", "Which claims are assumptions?", "Ledger now labels source type and confidence on every major claim."),
    ]
    table(slide, 0.72, 1.35, [1.35, 3.0, 6.65], 0.66,
          ["Persona", "Tough question", "Revision made"], rows,
          header_fill=COLORS["red"], font_size=7.7)
    callout(slide, 0.95, 5.7, 11.25, 0.6, "Why this matters",
            "The red-team loop is valuable only if it changes the output. Here it changed the ask, the approval conditions, and the confidence language.",
            accent=COLORS["red"], fill=COLORS["red_light"])
    note(slide, "This is the proof that the review loop is not ceremonial.")


def slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "quality comparison",
           "The orchestrated deck makes decision-readiness gates explicit")
    text(slide, 3.45, 1.32, 2.4, 0.22, "Generic prompt risk", size=8.8,
         bold=True, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    text(slide, 8.15, 1.32, 2.8, 0.22, "Orchestrated output", size=8.8,
         bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)
    rows = [
        ("Decision ask", "unclear or buried", "approve gated pilot / defer / reject rollout"),
        ("Framework use", "decorative model", "issue tree selected; BMC and Porter rejected"),
        ("Evidence discipline", "facts and assumptions mixed", "source type, confidence, gaps, appendix cue"),
        ("Alternatives", "single-path story", "scale now, gated pilot, defer compared"),
        ("Objections", "left for Q&A", "CFO/operator/sales/evidence objections handled"),
        ("Action", "next steps vague", "owner, proof gates, kill criteria, next review"),
    ]
    for i, (gate, before, after) in enumerate(rows):
        delta_row(slide, 1.72 + i * 0.68, gate, before, after, COLORS["teal"])
    callout(slide, 0.9, 6.08, 11.35, 0.48, "Proof standard",
         "This is a designed readiness comparison, not a claim that every generic AI output fails. A real pilot should test before/after quality on live decks.",
            accent=COLORS["amber"], fill=COLORS["amber_light"])
    note(slide, "Do not claim universal superiority. Say this is what the harness is designed to force.")


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "adoption guardrails",
           "Use the skill for decision decks; keep humans accountable for facts and judgment")
    columns = [
        ("Use it for", "strategy, GTM, product, board, investment, roadmap, workshop readout", COLORS["teal"]),
        ("Do not use it for", "pure formatting, unverified claims, confidential upload without approval, legal or financial certainty", COLORS["red"]),
        ("Human review required", "facts, source permission, assumptions, recommendations, risk tolerance, final business decision", COLORS["blue"]),
    ]
    for i, (title, body, color) in enumerate(columns):
        x = 0.75 + i * 4.0
        rect(slide, x, 1.45, 3.55, 1.42, fill=COLORS["white"], line=color)
        rect(slide, x, 1.45, 3.55, 0.14, fill=color, line=color)
        text(slide, x + 0.18, 1.82, 3.12, 0.18, title, size=9.4,
             bold=True, color=COLORS["ink"])
        text(slide, x + 0.18, 2.2, 3.12, 0.42, body, size=7.4,
             color=COLORS["muted"])
    gate_items = [
        ("Input", "audience, decision, evidence, time limit, constraints"),
        ("Evidence", "weak proof triggers research-needed or pilot-needed deck"),
        ("Decision", "output names approve, reject, pause, test, fund, or assign"),
        ("Owner", "a person remains accountable for truth and final recommendation"),
    ]
    rect(slide, 0.95, 3.65, 11.2, 1.12, fill=COLORS["white"], line=COLORS["line"])
    text(slide, 1.18, 3.98, 1.2, 0.16, "Gate checks", size=9.0,
         bold=True, color=COLORS["teal"])
    for i, (title, detail) in enumerate(gate_items):
        x = 2.55 + i * 2.35
        text(slide, x, 3.9, 1.85, 0.16, title, size=8.2, bold=True,
             color=COLORS["ink"], align=PP_ALIGN.CENTER)
        text(slide, x, 4.2, 1.85, 0.26, detail, size=6.6,
             color=COLORS["muted"], align=PP_ALIGN.CENTER)
        if i > 0:
            connector(slide, x - 0.25, 4.23, x - 0.02, 4.23, color=COLORS["line"], width=0.8)
    callout(slide, 0.95, 5.65, 11.25, 0.58, "Operating principle",
            "The skill can structure the work and expose gaps; it cannot manufacture proof, replace source review, or own the business decision.",
            accent=COLORS["amber"], fill=COLORS["amber_light"])
    note(slide, "This handles the skeptical executive's limitations question.")


def slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "verification",
           "The artifact path is reproducible and current package checks pass",
           "local repo validation on 2026-06-26")
    steps = [
        ("Build", "python-pptx wrote the PPTX artifact", COLORS["blue"]),
        ("Render", "LibreOffice exported a matching PDF", COLORS["teal"]),
        ("Inspect", "contact sheet reviewed after export", COLORS["green"]),
        ("Package", "release:check passed for 11 skills", COLORS["blue"]),
        ("Source", "reader signals feed research, not runtime routing", COLORS["amber"]),
    ]
    for i, (title, proof, color) in enumerate(steps, 1):
        x = 0.72 + (i - 1) * 2.38
        validation_step(slide, x, 1.75, i, title, proof, color)
        if i < len(steps):
            connector(slide, x + 2.1, 2.25, x + 2.3, 2.25, color=COLORS["line"], width=1.0)
    rect(slide, 1.0, 3.55, 11.15, 1.15, fill=COLORS["white"], line=COLORS["line"])
    text(slide, 1.22, 3.8, 2.05, 0.18, "What this proves", size=9.2,
         bold=True, color=COLORS["teal"])
    text(slide, 3.15, 3.78, 8.68, 0.22,
         "The deck artifact, source references, and public skill package are reproducible and locally validated.",
         size=9.2, bold=True, color=COLORS["ink"])
    text(slide, 3.15, 4.18, 8.68, 0.2,
         "It does not prove every future recommendation is true; facts, source permissions, and business judgment still require human review.",
         size=8.3, color=COLORS["muted"])
    callout(slide, 0.95, 5.65, 11.25, 0.56, "Scope of proof",
            "Technical validation proves package integrity and renderability. Advisory quality is validated only through evidence review and live before/after use.",
            accent=COLORS["amber"], fill=COLORS["amber_light"])
    note(slide, "This answers whether the skill is real and validated, while keeping proof claims scoped.")


def slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "decision close",
           "Decision readiness is the success measure, not slide count")
    rect(slide, 0.78, 1.45, 11.78, 1.15, fill=COLORS["teal"], line=COLORS["teal"])
    text(slide, 1.05, 1.72, 2.25, 0.18, "Final gate", size=9.0,
         bold=True, color=COLORS["white"])
    text(slide, 3.15, 1.66, 8.95, 0.32,
         "Use the skill when the cost of a weak deck is a weak decision.",
         size=17.0, bold=True, color=COLORS["white"])
    text(slide, 3.15, 2.12, 8.95, 0.18,
         "Unclear ask, late objections, hidden assumptions, or no accountable next owner.",
         size=8.4, color=COLORS["white"])
    gates = [
        ("1. Decision ask", "approve / reject / test is explicit", COLORS["blue"]),
        ("2. Evidence standard", "source type, confidence, and gaps are visible", COLORS["green"]),
        ("3. Human accountability", "owner, proof gates, and final judgment remain named", COLORS["amber"]),
    ]
    for i, (title, detail, color) in enumerate(gates):
        x = 0.92 + i * 4.02
        rect(slide, x, 3.18, 3.5, 1.28, fill=COLORS["white"], line=color)
        rect(slide, x, 3.18, 3.5, 0.12, fill=color, line=color)
        text(slide, x + 0.18, 3.48, 3.05, 0.18, title, size=9.1,
             bold=True, color=COLORS["ink"])
        text(slide, x + 0.18, 3.86, 3.05, 0.28, detail, size=7.7,
             color=COLORS["muted"])
    rect(slide, 1.12, 5.15, 10.95, 0.58, fill=COLORS["white"], line=COLORS["teal"])
    text(slide, 1.38, 5.36, 10.4, 0.14,
         "Pilot rubric: compare 3 real decks to baseline; require >=4/5 on ask, evidence, alternatives, owner, and rehearsal changes.",
         size=10.6, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)
    text(slide, 1.08, 6.16, 11.2, 0.24,
         "Ship as a public demo artifact; use live before/after evidence before claiming business-outcome improvement.",
         size=9.5, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)
    note(slide, "Close on decision readiness and a controlled next action.")


SLIDES = [
    slide_1,
    slide_2,
    slide_3,
    slide_4,
    slide_5,
    slide_6,
    slide_7,
    slide_8,
    slide_9,
    slide_10,
    slide_11,
    slide_12,
]


OUTLINE = [
    ("Use Presentation Orchestrator when AI slide generation must become decision preparation", "Recommendation and adoption ask."),
    ("The skill is worth using only when it can prove a stronger decision deck", "Executive summary with proof points."),
    ("A messy GTM request becomes a decision question the room can answer", "Synthetic mini case setup."),
    ("The skill chooses the decision job before it chooses the framework", "Selected and rejected methods with reasons."),
    ("The evidence ledger turns weak signals into explicit confidence and next checks", "Hero proof ledger."),
    ("The selected framework changes the decision: pilot beats scale-now and defer", "Issue tree plus option matrix."),
    ("The deck packet turns analysis into five action-title slides", "Storyboard output."),
    ("The rehearsal loop improves the recommendation before the room reacts", "Red-team objections and revisions."),
    ("The orchestrated deck makes decision-readiness gates explicit", "Before/after quality gates."),
    ("Use the skill for decision decks; keep humans accountable for facts and judgment", "Use / do not use / human review guardrails."),
    ("The artifact path is reproducible and current package checks pass", "Validation evidence."),
    ("Decision readiness is the success measure, not slide count", "Final readiness scorecard and next action."),
]


def build_markdown():
    lines = [
        "# Presentation Orchestrator: How the Skill Works",
        "",
        "This version was rebuilt after a consultation-team review. It is case-led, answer-first, and uses decision-readiness as the quality bar.",
        "",
        "## Mainline Storyline",
        "",
    ]
    for i, (title, point) in enumerate(OUTLINE, 1):
        lines.append(f"{i}. **{title}** - {point}")
    lines.extend([
        "",
        "## Quality Bar Applied",
        "",
        "- Lead with recommendation and decision ask.",
        "- Prove the skill through one synthetic but realistic mini case.",
        "- Show selected and rejected frameworks with reasons.",
        "- Make evidence gaps, confidence, and red-team revisions visible.",
        "- Remove internal implementation detail from the mainline story.",
        "- Use reader signals to source and validate diagram-family candidates; route runtime choices by decision job and evidence fit.",
        "",
        "## Key Source Anchors",
        "",
        "- `skills/presentation-orchestrator/SKILL.md`",
        "- `skills/presentation-orchestrator/references/methodology-router.md`",
        "- `skills/presentation-orchestrator/references/presentation-diagram-source-catalog.md`",
        "- `skills/presentation-orchestrator/references/deck-output-template.md`",
        "- `skills/presentation-orchestrator/references/harness-loop-contract.md`",
        "- `docs/research/presentation-diagram-source-research.md`",
        "",
    ])
    MD_PATH.write_text("\n".join(lines), encoding="utf-8")


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for builder in SLIDES:
        builder(prs)
    prs.save(PPTX_PATH)
    build_markdown()
    print(PPTX_PATH.relative_to(ROOT))
    print(MD_PATH.relative_to(ROOT))


if __name__ == "__main__":
    main()
